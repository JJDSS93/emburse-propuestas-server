from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io, json

app = Flask(__name__)
CORS(app)

def hex_color(h):
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

AZUL   = hex_color('0D2240')
AZUL2  = hex_color('1E6DB5')
BLANCO = hex_color('FFFFFF')
GRIS   = hex_color('666666')
PAR    = hex_color('EDF2F7')
IMPAR  = hex_color('FFFFFF')

ROW_H = Cm(1.0)   # altura fija por fila: 1 cm
HEAD_H = Cm(1.2)  # altura cabecera: 1.2 cm

def clear_slide(slide):
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape._element)

def add_textbox(slide, text, x, y, w, h, bold=False, size=32, color=None, align=PP_ALIGN.LEFT):
    if color is None: color = AZUL
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return tb

def add_table(slide, headers, rows, x, y, w, col_widths):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    # Altura dinámica: cabecera + filas de datos
    h = HEAD_H + ROW_H * len(rows)
    tbl = slide.shapes.add_table(n_rows, n_cols, Emu(x), Emu(y), Emu(w), h).table
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = Emu(cw)
    # Altura de cada fila
    tbl.rows[0].height = HEAD_H
    for i in range(1, n_rows):
        tbl.rows[i].height = ROW_H
    # Cabecera
    for j, ht in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = ht
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZUL
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
        run.font.bold = True
        run.font.color.rgb = BLANCO
        run.font.size = Pt(11)
        run.font.name = 'Calibri'
    # Filas de datos
    for i, row in enumerate(rows):
        bg = PAR if i % 2 == 0 else IMPAR
        for j, val in enumerate(row):
            cell = tbl.cell(i+1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            if p.runs:
                run = p.runs[0]
            else:
                run = p.add_run()
            run.font.color.rgb = AZUL
            run.font.size = Pt(11)
            run.font.name = 'Calibri'
    return h

@app.route('/health')
def health():
    return jsonify({'status': 'ok'})

@app.route('/generar', methods=['POST', 'OPTIONS'])
def generar():
    if request.method == 'OPTIONS':
        resp = app.make_default_options_response()
        resp.headers['Access-Control-Allow-Origin'] = '*'
        resp.headers['Access-Control-Allow-Headers'] = 'Content-Type'
        resp.headers['Access-Control-Allow-Methods'] = 'POST, OPTIONS'
        return resp

    try:
        pptx_file = request.files.get('plantilla')
        if not pptx_file:
            return jsonify({'error': 'No se recibió la plantilla'}), 400

        datos = json.loads(request.form.get('datos', '{}'))
        empresa       = datos.get('empresa', 'Cliente')
        moneda        = datos.get('moneda', 'USD')
        licencias     = datos.get('licencias', [])
        setup         = datos.get('setup', [])
        opcionales    = datos.get('opcionales', [])
        notas         = datos.get('notas', 'Precios en USD. IVA no incluido. Vigencia 30 días.')
        fecha         = datos.get('fecha', '')
        leyenda_lic   = (datos.get('leyenda_lic') or '').strip()
        leyenda_setup = (datos.get('leyenda_setup') or '').strip()
        leyenda_opc   = (datos.get('leyenda_opc') or '').strip()

        def fmt(n):
            sym = '$' if moneda in ('USD','MXN','COP') else '€'
            return f"{sym}{float(n):,.2f} {moneda}"

        prs = Presentation(io.BytesIO(pptx_file.read()))

        # Logo en slide 1 — centrado horizontalmente, bajado 2cm desde el centro vertical
        logo_file = request.files.get('logo')
        if logo_file:
            logo_bytes = io.BytesIO(logo_file.read())
            slide1 = prs.slides[0]
            slide_w = prs.slide_width
            slide_h = prs.slide_height
            logo_w = Cm(5)
            logo_h = Cm(5)
            logo_x = (slide_w - logo_w) // 2
            logo_y = (slide_h - logo_h) // 2 + Cm(1.5)
            slide1.shapes.add_picture(logo_bytes, logo_x, logo_y, logo_w, logo_h)

        # Preparar filas — si no hay datos, añadir una fila vacía para que la tabla no quede rara
        lic_rows = [[r.get('nombre',''), str(r.get('vol','')),
            fmt(r['precio']) if r.get('precio') else '',
            fmt(float(r.get('vol') or 0)*float(r.get('precio') or 0)) if r.get('nombre') else '']
            for r in licencias] or [['', '', '', '']]

        setup_rows = [[r.get('servicio',''), r.get('descripcion',''), r.get('tipo','Único'),
            fmt(r['precio']) if r.get('precio') else ''] for r in setup] or [['', '', '', '']]

        opc_rows = [[r.get('servicio',''), r.get('descripcion',''), r.get('tipo','Anual'),
            fmt(r['precio']) if r.get('precio') else ''] for r in opcionales] or [['', '', '', '']]

        lic_total   = sum(float(r.get('vol') or 0)*float(r.get('precio') or 0) for r in licencias)
        setup_total = sum(float(r.get('precio') or 0) for r in setup)
        opc_total   = sum(float(r.get('precio') or 0) for r in opcionales)
        res_rows = [
            ["Pago anual recurrente", "Licencias únicamente", fmt(lic_total)],
            ["Pago anual más servicios (sólo primer año)", "Licencias + Setup & Onboarding", fmt(lic_total+setup_total)],
            ["Pago anual más servicios más opcionales", "Licencias + Setup + Servicios Opcionales", fmt(lic_total+setup_total+opc_total)],
        ]

        Y_TITULO   = 304800
        Y_TABLA    = 990600
        X          = 457200
        W          = 8229600
        Y_LEYENDA_OFFSET = Cm(0.4)  # espacio entre tabla y leyenda

        # Slide 19 — Licencias
        s19 = prs.slides[18]
        clear_slide(s19)
        add_textbox(s19, "Licencias", X, Y_TITULO, W, 533400, bold=True, size=32)
        h19 = add_table(s19, ["Producto / Licencia","Volumen",f"Precio unit. ({moneda})",f"Total ({moneda})"],
                  lic_rows, X, Y_TABLA, W, [2500000,1900000,1900000,1700000])
        if leyenda_lic:
            add_textbox(s19, leyenda_lic, X, Y_TABLA + h19 + Y_LEYENDA_OFFSET, W, 304800, size=10, color=GRIS)

        # Slide 20 — Setup
        s20 = prs.slides[19]
        clear_slide(s20)
        add_textbox(s20, "Setup & Onboarding", X, Y_TITULO, W, 533400, bold=True, size=32)
        h20 = add_table(s20, ["Servicio","Descripción","Tipo de pago",f"Precio ({moneda})"],
                  setup_rows, X, Y_TABLA, W, [2000000,3100000,1500000,1629600])
        if leyenda_setup:
            add_textbox(s20, leyenda_setup, X, Y_TABLA + h20 + Y_LEYENDA_OFFSET, W, 304800, size=10, color=GRIS)

        # Slide 21 — Opcionales
        s21 = prs.slides[20]
        clear_slide(s21)
        add_textbox(s21, "Servicios Opcionales", X, Y_TITULO, W, 533400, bold=True, size=32)
        h21 = add_table(s21, ["Servicio","Descripción","Tipo de pago",f"Precio ({moneda})"],
                  opc_rows, X, Y_TABLA, W, [2000000,3100000,1500000,1629600])
        if leyenda_opc:
            add_textbox(s21, leyenda_opc, X, Y_TABLA + h21 + Y_LEYENDA_OFFSET, W, 304800, size=10, color=GRIS)

        # Slide 22 — Resumen
        s22 = prs.slides[21]
        clear_slide(s22)
        add_textbox(s22, "Resumen", X, Y_TITULO, W, 533400, bold=True, size=32)
        h22 = add_table(s22, ["Escenario","Descripción",f"Inversión Año 1 ({moneda})"],
                  res_rows, X, Y_TABLA, W, [2300000,3800000,2129600])
        add_textbox(s22, "Condiciones y notas adicionales:", X, Y_TABLA + h22 + Y_LEYENDA_OFFSET, W, 304800, bold=True, size=14, color=AZUL2)
        notas_text = "\n".join(f"• {l.strip()}" for l in notas.split('\n') if l.strip())
        add_textbox(s22, notas_text, X, Y_TABLA + h22 + Y_LEYENDA_OFFSET + Cm(0.8), W, 700000, size=10, color=GRIS)

        out = io.BytesIO()
        prs.save(out)
        out.seek(0)

        fecha_str = (fecha or '').replace('-','')
        filename = f"Propuesta_Emburse_{empresa.replace(' ','_')}_{fecha_str}.pptx"

        resp = send_file(out,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True, download_name=filename)
        resp.headers['Access-Control-Allow-Origin'] = '*'
        return resp

    except Exception as e:
        import traceback
        print(traceback.format_exc())
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    import os
    port = int(os.environ.get('PORT', 8080))
    app.run(host='0.0.0.0', port=port)
